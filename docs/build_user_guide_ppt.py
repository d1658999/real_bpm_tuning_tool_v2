r"""Build the Real BPM Tuning Tool user-guide PowerPoint.

The script also captures two current GUI screenshots with a non-confidential
demo Touchstone file so the deck stays aligned with the application layout.

Run from the repository root:

    .\.venv\Scripts\python.exe docs\build_user_guide_ppt.py
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Iterable, Sequence

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
os.environ.setdefault("QT_QPA_FONTDIR", str(Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets" / "user-guide"
OUTPUT = DOCS / "Real_BPM_Tuning_Tool_User_Guide.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F5F5F7"
WHITE = "FFFFFF"
INK = "1D1D1F"
SECONDARY = "6E6E73"
BLUE = "0066CC"
CYAN = "5AC8FA"
GREEN = "248A3D"
ORANGE = "D97706"
RED = "D70015"
PALE_BLUE = "EAF3FC"
PALE_GREEN = "EAF7ED"
PALE_ORANGE = "FFF4E5"
PALE_RED = "FDECEC"
LINE = "D2D2D7"
NAVY = "0B1F33"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str | None = None,
    radius: bool = True,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    segments: Sequence[tuple[str, bool, str]],
    *,
    size: float = 18,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text_value, bold, color in segments:
        run = paragraph.add_run()
        run.text = text_value
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Iterable[str],
    *,
    size: float = 17,
    color: str = INK,
    accent: str = BLUE,
    spacing: float = 8,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        paragraph.level = 0
        paragraph.text = ""
        marker = paragraph.add_run()
        marker.text = "●  "
        marker.font.name = "Aptos"
        marker.font.size = Pt(size - 2)
        marker.font.color.rgb = rgb(accent)
        run = paragraph.add_run()
        run.text = item
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, 0.55, 0.32, 12.2, 0.48, title, size=27, bold=True, color=INK)
    if subtitle:
        add_text(slide, 0.57, 0.83, 12.0, 0.32, subtitle, size=11.5, color=SECONDARY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.56), Inches(1.16), Inches(0.62), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(BLUE)
    line.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, 0.56, 7.12, 5.0, 0.2, "Real BPM Tuning Tool  •  User guide", size=8.5, color=SECONDARY)
    add_text(slide, 12.1, 7.08, 0.66, 0.23, str(number), size=9, color=SECONDARY, align=PP_ALIGN.RIGHT)


def add_badge(slide, x: float, y: float, label: str, *, fill: str = BLUE, size: float = 0.34) -> None:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(fill)
    circle.line.fill.background()
    add_text(
        slide,
        x,
        y + 0.002,
        size,
        size - 0.004,
        label,
        size=12,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_label_pill(slide, x: float, y: float, w: float, text: str, *, fill: str, color: str = WHITE) -> None:
    add_rect(slide, x, y, w, 0.34, fill=fill, radius=True)
    add_text(slide, x, y + 0.015, w, 0.27, text, size=10.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    draw_w = image_w * scale
    draw_h = image_h * scale
    return slide.shapes.add_picture(
        str(path),
        Inches(x + (w - draw_w) / 2),
        Inches(y + (h - draw_h) / 2),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )


def capture_gui_assets() -> tuple[Path, Path]:
    """Capture configured and cascaded GUI states using a public demo fixture."""

    from PyQt5.QtGui import QFont, QFontDatabase
    from PyQt5.QtWidgets import QApplication, QCheckBox, QDoubleSpinBox

    from bpm_tuner.gui import CoreAPI, MainWindow

    ASSETS.mkdir(parents=True, exist_ok=True)
    configured_path = ASSETS / "bpm_tuner_configured.png"
    cascaded_path = ASSETS / "bpm_tuner_cascade_result.png"

    app = QApplication.instance() or QApplication([])
    fonts_dir = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    for font_name in ("segoeui.ttf", "segoeuib.ttf"):
        font_path = fonts_dir / font_name
        if font_path.is_file():
            QFontDatabase.addApplicationFont(str(font_path))
    app.setFont(QFont("Segoe UI", 9))
    window = MainWindow(ROOT)
    window.resize(1600, 920)
    demo_path = ROOT / "snp_files" / "Cubs_RJF_DK2p55.s4p"
    window.files_panel.set_paths([demo_path])
    window.candidate_count.setValue(2)

    def set_mode(row: int, mode: str, signal: str | None = None) -> None:
        controls = window.port_panel.row_controls(row)
        controls["mode"].setCurrentText(mode)
        if signal:
            controls["signal"].setCurrentText(signal)

    set_mode(0, "signal", "s1")
    set_mode(1, "signal", "s2")
    set_mode(2, "open/inductor/capacitor")
    set_mode(3, "open")
    controls = window.port_panel.row_controls(2)
    controls["inductor_min_nh"].setValue(0.2)
    controls["inductor_max_nh"].setValue(3.0)
    controls["capacitor_min_pf"].setValue(0.2)
    controls["capacitor_max_pf"].setValue(10.0)
    window.port_panel._refresh_target_controls()

    for column, value in ((0, 3.3), (1, 5.0)):
        spin = window.port_panel.frequency_target_table.cellWidget(0, column)
        if isinstance(spin, QDoubleSpinBox):
            spin.setValue(value)
    enabled = window.port_panel.frequency_target_table.cellWidget(0, 2)
    if isinstance(enabled, QCheckBox):
        enabled.setChecked(True)

    window.show()
    app.processEvents()
    window.grab().save(str(configured_path))

    try:
        result = CoreAPI.run_cascade(window.current_config())
        window._set_result(result)
        window.status_label.setText("Cascade complete")
    except Exception as exc:  # Keep deck generation useful even if RF deps drift.
        window.status_label.setText(f"Cascade preview unavailable: {exc}")
    app.processEvents()
    window.grab().save(str(cascaded_path))
    window.close()
    app.processEvents()
    return configured_path, cascaded_path


def new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    return presentation


def add_slide(presentation: Presentation, number: int):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(BG)
    add_footer(slide, number)
    return slide


def slide_cover(presentation: Presentation, gui_result: Path) -> None:
    slide = add_slide(presentation, 1)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=NAVY, radius=False)
    add_label_pill(slide, 0.72, 0.66, 1.45, "USER GUIDE", fill=BLUE)
    add_text(slide, 0.72, 1.28, 6.1, 0.58, "Real BPM", size=38, color=WHITE, bold=True)
    add_text(slide, 0.72, 1.9, 6.1, 0.58, "Tuning Tool", size=38, color=WHITE, bold=True)
    add_text(
        slide,
        0.75,
        3.12,
        5.4,
        0.8,
        "Configure multiport Touchstone networks, evaluate RF performance, and select measured BOM parts with production-aware optimization.",
        size=17.5,
        color="D6E6F5",
    )
    add_text(slide, 0.75, 5.75, 4.9, 0.44, "Python / PyQt interface  •  Rust exhaustive sweep", size=12, color="9FC7E8")
    add_rect(slide, 6.55, 0.66, 6.1, 5.98, fill=WHITE, radius=True)
    add_picture_contain(slide, gui_result, 6.72, 0.84, 5.76, 5.62)
    add_text(slide, 7.05, 6.73, 5.1, 0.25, "Example cascade view using a non-confidential demo file", size=9, color="B8CADB", align=PP_ALIGN.CENTER)


def slide_workflow(presentation: Presentation) -> None:
    slide = add_slide(presentation, 2)
    add_title(slide, "The six-step workflow", "Use Cascade to validate the configured circuit before committing time to Optimization.")
    steps = [
        ("1", "Add", "Touchstone files"),
        ("2", "Configure", "every active port"),
        ("3", "Set", "bands and targets"),
        ("4", "Cascade", "check the circuit"),
        ("5", "Optimize", "compare five strategies"),
        ("6", "Export", "config, SNP, CSV, plots"),
    ]
    x_positions = [0.65, 2.75, 4.85, 6.95, 9.05, 11.15]
    for index, ((number, verb, noun), x) in enumerate(zip(steps, x_positions)):
        if index < len(steps) - 1:
            connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.61), Inches(3.18), Inches(0.55), Inches(0.045))
            connector.fill.solid()
            connector.fill.fore_color.rgb = rgb(LINE)
            connector.line.fill.background()
        add_rect(slide, x, 2.15, 1.55, 2.18, fill=WHITE, line=LINE)
        add_badge(slide, x + 0.56, 2.43, number, size=0.43)
        add_text(slide, x + 0.12, 3.08, 1.31, 0.35, verb, size=18, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.14, 3.51, 1.27, 0.52, noun, size=11.2, color=SECONDARY, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.08, 5.15, 11.18, 0.88, fill=PALE_BLUE, line="C9DFF3")
    add_rich_text(
        slide,
        1.35,
        5.43,
        10.7,
        0.32,
        [
            ("Fast feedback: ", True, BLUE),
            ("Run Cascade after any connection, frequency, or component change. Export buttons become available only after a successful result.", False, INK),
        ],
        size=14,
        align=PP_ALIGN.CENTER,
    )


def slide_launch(presentation: Presentation) -> None:
    slide = add_slide(presentation, 3)
    add_title(slide, "Launch the tool and keep its measured BOM beside it", "The Windows EXE needs no Python or Rust installation on the target PC.")
    add_rect(slide, 0.6, 1.52, 5.95, 4.75, fill=WHITE, line=LINE)
    add_label_pill(slide, 0.9, 1.82, 1.9, "DEPLOYED EXE", fill=BLUE)
    add_text(slide, 0.92, 2.42, 4.9, 0.44, "Double-click BPMTuningTool.exe", size=21, bold=True)
    add_text(slide, 0.92, 3.05, 4.8, 1.55, "dist\\\n  BPMTuningTool.exe\n  Capacitors_BOM\\\n  Inductors_BOM\\", size=17, color=NAVY, font="Consolas")
    add_rect(slide, 0.92, 5.02, 5.05, 0.74, fill=PALE_ORANGE, line="F5D4A5")
    add_text(slide, 1.12, 5.22, 4.65, 0.34, "Keep both BOM folder names unchanged and move all three items together.", size=12.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 6.82, 1.52, 5.91, 4.75, fill=WHITE, line=LINE)
    add_label_pill(slide, 7.12, 1.82, 1.7, "DEVELOPER", fill=NAVY)
    add_text(slide, 7.12, 2.42, 4.9, 0.44, "Launch from PowerShell", size=21, bold=True)
    add_rect(slide, 7.12, 3.02, 5.0, 0.65, fill=NAVY, radius=True)
    add_text(slide, 7.34, 3.22, 4.55, 0.28, "bpm-tuner --gui", size=16, color=WHITE, font="Consolas")
    add_bullets(
        slide,
        7.15,
        4.08,
        5.0,
        1.35,
        ["Start in the project root.", "For a first optimization, set BOM samples/type to 2."],
        size=14.5,
    )
    add_text(slide, 7.15, 5.54, 4.95, 0.34, "Build the deployable EXE with .\\build_one_exe.bat", size=12, color=SECONDARY)


def slide_interface(presentation: Presentation, gui_configured: Path) -> None:
    slide = add_slide(presentation, 4)
    add_title(slide, "Interface tour", "Work left to right: source files → port and target settings → RF plots.")
    add_rect(slide, 0.46, 1.38, 12.42, 5.28, fill=WHITE, line=LINE)
    add_picture_contain(slide, gui_configured, 0.58, 1.52, 12.18, 4.18)
    positions = [(0.86, 2.12, "1"), (3.23, 2.12, "2"), (7.9, 2.12, "3"), (1.45, 1.63, "4")]
    for x, y, label in positions:
        add_badge(slide, x, y, label, size=0.36)
    captions = [
        ("1", "Files", "Add or remove .sNp networks."),
        ("2", "Configuration", "Assign modes, signals, bands, targets, and ranges."),
        ("3", "Plots", "Smith, S21, VSWR, and return loss."),
        ("4", "Toolbar", "Run, save/load, and export."),
    ]
    x_values = [0.7, 3.78, 6.86, 9.94]
    for x, (number, label, detail) in zip(x_values, captions):
        add_rich_text(slide, x, 5.9, 2.78, 0.26, [(number + "  ", True, BLUE), (label, True, INK)], size=11.5)
        add_text(slide, x + 0.26, 6.17, 2.54, 0.32, detail, size=9.6, color=SECONDARY)


def slide_ports(presentation: Presentation) -> None:
    slide = add_slide(presentation, 5)
    add_title(slide, "Configure every active port", "The selected Mode controls what appears in the Port configuration column.")
    modes = [
        ("signal", "External driven/dependent port", PALE_BLUE, BLUE),
        ("connect", "Reciprocal link to another file/port", "F0F0F2", INK),
        ("open / short", "Unused port or 0 Ω termination", "F0F0F2", SECONDARY),
        ("inductor / capacitor", "Measured part or optimized range", PALE_GREEN, GREEN),
        ("inductor/capacitor", "Search both measured catalogs", PALE_GREEN, GREEN),
        ("open/inductor/capacitor", "True open baseline + both catalogs", PALE_ORANGE, ORANGE),
    ]
    for index, (label, detail, fill, accent) in enumerate(modes):
        row, column = divmod(index, 2)
        x = 0.62 + column * 6.18
        y = 1.48 + row * 1.22
        add_rect(slide, x, y, 5.86, 0.98, fill=fill, line=LINE)
        add_text(slide, x + 0.25, y + 0.18, 2.44, 0.32, label, size=15, bold=True, color=accent)
        add_text(slide, x + 2.72, y + 0.18, 2.9, 0.55, detail, size=12.2, color=INK)
    add_rect(slide, 0.62, 5.28, 12.04, 1.15, fill=WHITE, line=LINE)
    rules = [
        "Connections must be reciprocal.",
        "Signals must be unique and consecutive from s1; use 2–4 total.",
        "The highest-numbered signal is the dependent antenna/common port.",
        "A file not used yet may remain loaded only when all of its ports are open.",
    ]
    add_text(slide, 0.9, 5.56, 2.05, 0.33, "Preflight rules", size=17, bold=True, color=RED)
    add_text(slide, 2.95, 5.44, 9.34, 0.69, "  •  ".join(rules), size=11.5, color=INK, valign=MSO_ANCHOR.MIDDLE)


def slide_targets(presentation: Presentation) -> None:
    slide = add_slide(presentation, 6)
    add_title(slide, "Set frequency bands and optional Smith targets", "Only driven signals appear in the compact target table; the dependent final signal is intentionally hidden.")
    add_rect(slide, 0.62, 1.48, 12.05, 1.34, fill=WHITE, line=LINE)
    headers = ["Signal", "Start GHz", "Stop GHz", "Target", "Target R Ω", "Target X Ω"]
    values = ["s1", "3.3", "5.0", "☑ Enable", "50.0000", "0.0000"]
    widths = [0.9, 1.75, 1.75, 1.8, 2.25, 2.25]
    x = 0.91
    for header, value, width in zip(headers, values, widths):
        add_text(slide, x, 1.73, width, 0.25, header, size=10, color=SECONDARY, bold=True)
        add_text(slide, x, 2.13, width, 0.34, value, size=14, color=INK, bold=header == "Signal")
        x += width
    add_rect(slide, 0.62, 3.22, 5.84, 2.65, fill=PALE_BLUE, line="C9DFF3")
    add_text(slide, 0.92, 3.55, 4.9, 0.38, "Frequency", size=20, bold=True, color=BLUE)
    add_bullets(
        slide,
        0.92,
        4.08,
        4.98,
        1.35,
        [
            "Auto uses the full Touchstone overlap.",
            "Enter both limits or leave both Auto.",
            "Start must be lower than Stop.",
        ],
        size=14.5,
    )
    add_rect(slide, 6.82, 3.22, 5.85, 2.65, fill=PALE_GREEN, line="C6E6CD")
    add_text(slide, 7.12, 3.55, 4.9, 0.38, "Smith target", size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        7.12,
        4.08,
        4.98,
        1.35,
        [
            "Disabled by default.",
            "Enter physical R + jX in ohms.",
            "50 + j0 Ω is the 50 Ω Smith-chart center.",
        ],
        size=14.5,
        accent=GREEN,
    )
    add_text(slide, 0.78, 6.26, 11.8, 0.37, "Tip: set different bands and targets for s1, s2, and s3 when they are driven ports; never assign an independent target to the dependent final signal.", size=12.5, color=SECONDARY, align=PP_ALIGN.CENTER)


def slide_search_size(presentation: Presentation) -> None:
    slide = add_slide(presentation, 7)
    add_title(slide, "Control BOM coverage without exploding the search", "BOM samples/type is global; each tunable port may have its own inclusive L and C range.")
    add_rect(slide, 0.62, 1.5, 5.8, 4.97, fill=WHITE, line=LINE)
    add_text(slide, 0.94, 1.82, 4.9, 0.42, "Choices per tunable port", size=21, bold=True)
    rows = [
        ("inductor", "N"),
        ("capacitor", "N"),
        ("inductor/capacitor", "2N"),
        ("open/inductor/capacitor", "2N + 1"),
    ]
    for index, (mode, count) in enumerate(rows):
        y = 2.54 + index * 0.66
        add_text(slide, 1.02, y, 3.55, 0.3, mode, size=14.5, color=INK)
        add_label_pill(slide, 4.83, y - 0.04, 0.92, count, fill=BLUE)
    add_text(slide, 1.02, 5.42, 4.92, 0.66, "Total combinations = product of the choices at every tunable port", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, 6.78, 1.5, 5.89, 4.97, fill=NAVY, line=NAVY)
    add_text(slide, 7.12, 1.83, 5.0, 0.38, "Safe first run", size=21, bold=True, color=WHITE)
    add_text(slide, 7.12, 2.52, 4.95, 0.52, "Set BOM samples/type to 2", size=26, bold=True, color=CYAN)
    add_text(slide, 7.12, 3.19, 4.92, 0.54, "Example: three flexible open/L/C ports", size=14.5, color="D6E6F5")
    add_text(slide, 7.12, 3.83, 4.95, 0.65, "(2×2 + 1)³ = 5³ = 125", size=27, bold=True, color=WHITE, font="Consolas", align=PP_ALIGN.CENTER)
    add_text(slide, 7.12, 4.74, 5.0, 0.95, "Increase N only after the topology and bands are correct. The GUI warns above 100,000,000 combinations or an estimated 10-minute native sweep.", size=13.5, color="D6E6F5", align=PP_ALIGN.CENTER)


def slide_cascade(presentation: Presentation, gui_result: Path) -> None:
    slide = add_slide(presentation, 8)
    add_title(slide, "Run Cascade first, then read all four views", "A successful cascade enables SNP and insertion-loss CSV export.")
    add_rect(slide, 0.46, 1.36, 9.08, 5.58, fill=WHITE, line=LINE)
    add_picture_contain(slide, gui_result, 0.58, 1.49, 8.84, 5.24)
    labels = [
        (9.79, 1.55, "Smith chart", "S11/S22 location, contour, and target proximity"),
        (9.79, 2.73, "S21", "Insertion loss versus frequency"),
        (9.79, 3.91, "VSWR", "Mismatch severity; lower is better"),
        (9.79, 5.09, "Return loss", "Reflection magnitude in dB; higher is better"),
    ]
    for index, (x, y, title, detail) in enumerate(labels, start=1):
        add_badge(slide, x, y, str(index), size=0.34)
        add_text(slide, x + 0.48, y - 0.01, 2.48, 0.28, title, size=14.5, bold=True, color=BLUE)
        add_text(slide, x + 0.48, y + 0.34, 2.44, 0.58, detail, size=10.5, color=SECONDARY)
    add_text(slide, 9.82, 6.34, 2.73, 0.38, "Marker mode: click a trace to synchronize the readout across plots.", size=10.5, color=INK, bold=True)


def slide_optimization(presentation: Presentation) -> None:
    slide = add_slide(presentation, 9)
    add_title(slide, "Run Optimization to compare five engineering priorities", "Rust evaluates the full sampled Cartesian sweep; every strategy selects from the same candidate set.")
    strategies = [
        ("1", "minimum_bom", "Fewer selected components", BLUE),
        ("2", "balanced", "Balanced RF and production metrics", GREEN),
        ("3", "minimum_target", "Lowest Smith-target error", ORANGE),
        ("4", "smith_contour", "Compact impedance contour", "7C3AED"),
        ("5", "minimum_insertion_loss", "Lowest worst S21 loss", RED),
    ]
    for index, (number, key, detail, accent) in enumerate(strategies):
        x = 0.6 + index * 2.52
        add_rect(slide, x, 1.64, 2.22, 2.22, fill=WHITE, line=LINE)
        add_badge(slide, x + 0.18, 1.85, number, fill=accent, size=0.4)
        add_text(slide, x + 0.18, 2.46, 1.86, 0.58, key.replace("_", "\n"), size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 3.2, 1.86, 0.42, detail, size=10.3, color=SECONDARY, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.6, 4.32, 12.06, 1.72, fill=NAVY, line=NAVY)
    add_text(slide, 0.94, 4.61, 2.27, 0.4, "While it runs", size=20, bold=True, color=WHITE)
    add_bullets(
        slide,
        3.12,
        4.55,
        4.15,
        1.05,
        ["Watch progress and status.", "Cancel if the search is too large."],
        size=13.5,
        color=WHITE,
        accent=CYAN,
    )
    add_text(slide, 7.72, 4.59, 4.4, 0.68, "When complete, the Principal Engineer selects the lowest normalized production-risk score.", size=15, color="D6E6F5", align=PP_ALIGN.CENTER)
    add_text(slide, 1.02, 6.38, 11.32, 0.35, "Each winner is re-evaluated with independent 1.00 / 0.95 / 1.05 component-value factors before risk scoring.", size=12.5, color=SECONDARY, align=PP_ALIGN.CENTER)


def slide_results(presentation: Presentation, decision_plot: Path) -> None:
    slide = add_slide(presentation, 10)
    add_title(slide, "Review the fleet winner—and the trade-offs", "Example output from a repo-local optimization run; values will differ with your circuit and settings.")
    add_rect(slide, 0.48, 1.37, 8.25, 5.38, fill=WHITE, line=LINE)
    add_picture_contain(slide, decision_plot, 0.62, 1.51, 7.97, 5.1)
    add_rect(slide, 8.99, 1.37, 3.88, 2.72, fill=WHITE, line=LINE)
    add_text(slide, 9.28, 1.7, 3.18, 0.38, "Check these metrics", size=19, bold=True)
    add_bullets(
        slide,
        9.28,
        2.22,
        3.16,
        1.45,
        ["VSWR S11 / S22", "Worst insertion loss", "Target error", "Component count", "±5% sensitivity and risk"],
        size=12.3,
        spacing=4,
    )
    add_rect(slide, 8.99, 4.31, 3.88, 2.44, fill=PALE_ORANGE, line="F5D4A5")
    add_text(slide, 9.28, 4.62, 3.18, 0.35, "Production caveat", size=18, bold=True, color=ORANGE)
    add_text(slide, 9.28, 5.12, 3.17, 1.1, "The ±5% study is an electrical sensitivity proxy—not measured yield, vendor tolerance, or process capability. Validate PCB, fixture, temperature, bias, and production variation before release.", size=11.5, color=INK)


def slide_exports(presentation: Presentation) -> None:
    slide = add_slide(presentation, 11)
    add_title(slide, "Save, share, and export", "Save the project configuration separately from RF results so another user can reproduce the setup.")
    cards = [
        ("Save Config", "JSON", "Store files, port modes, targets, and per-port L/C ranges.", BLUE),
        ("Load Config", "JSON", "Restore the configuration; confirm paths resolve on the new PC.", NAVY),
        ("Export SNP", ".sNp", "Write the current cascaded or optimized Touchstone network.", GREEN),
        ("Export IL CSV", ".csv", "Write the current insertion-loss trace for analysis or reporting.", ORANGE),
    ]
    for index, (title, kind, detail, accent) in enumerate(cards):
        row, column = divmod(index, 2)
        x = 0.62 + column * 6.18
        y = 1.52 + row * 1.63
        add_rect(slide, x, y, 5.87, 1.34, fill=WHITE, line=LINE)
        add_label_pill(slide, x + 0.25, y + 0.24, 1.17, kind, fill=accent)
        add_text(slide, x + 1.68, y + 0.2, 3.75, 0.34, title, size=18, bold=True, color=accent)
        add_text(slide, x + 1.68, y + 0.63, 3.75, 0.45, detail, size=11.5, color=SECONDARY)
    add_rect(slide, 0.62, 5.02, 12.05, 1.27, fill=NAVY, line=NAVY)
    add_text(slide, 0.96, 5.31, 3.1, 0.36, "Optimization result folder", size=20, bold=True, color=WHITE)
    add_text(slide, 4.23, 5.18, 7.98, 0.72, "Fleet_results_YYYYMMDD_HHMMSS  →  five agent JSON files + five PNGs + agent_comparison.png + final_decision.png + report.md", size=13.5, color="D6E6F5", font="Consolas", valign=MSO_ANCHOR.MIDDLE)


def slide_troubleshooting(presentation: Presentation) -> None:
    slide = add_slide(presentation, 12)
    add_title(slide, "Preflight checklist and common warnings", "Most failures are actionable configuration issues; fix the named file and port, then run Cascade again.")
    checks = [
        ("No files", "Add at least one valid .sNp file; two or more are typical for a cascade."),
        ("Connection target missing", "Choose a destination and make the reverse port point back."),
        ("Signal assignment invalid", "Use unique, consecutive names from s1 with 2–4 total signals."),
        ("Frequency range invalid", "Enter both limits, keep Start < Stop, and stay inside the common data range."),
        ("Component range invalid", "Use positive Min ≤ Max and include at least one measured BOM part."),
        ("Search too large", "Reduce BOM samples/type or reduce the number of tunable ports."),
    ]
    for index, (problem, action) in enumerate(checks):
        row, column = divmod(index, 2)
        x = 0.62 + column * 6.18
        y = 1.43 + row * 1.39
        add_rect(slide, x, y, 5.88, 1.08, fill=WHITE, line=LINE)
        add_badge(slide, x + 0.22, y + 0.21, "!", fill=RED, size=0.38)
        add_text(slide, x + 0.77, y + 0.17, 1.92, 0.62, problem, size=14.5, bold=True, color=RED)
        add_text(slide, x + 2.92, y + 0.16, 2.64, 0.7, action, size=10.8, color=INK)
    add_rect(slide, 0.62, 5.84, 12.04, 0.72, fill=PALE_GREEN, line="C6E6CD")
    add_text(slide, 0.92, 6.06, 11.42, 0.27, "Ready to optimize when Cascade succeeds, the plotted band is correct, targets are intentional, and the estimated search size is acceptable.", size=13.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


def build() -> Path:
    gui_configured, gui_result = capture_gui_assets()
    decision_plot = ROOT / "Fleet_results_20260718_183635" / "final_decision.png"
    if not decision_plot.is_file():
        decision_plot = gui_result

    presentation = new_presentation()
    slide_cover(presentation, gui_result)
    slide_workflow(presentation)
    slide_launch(presentation)
    slide_interface(presentation, gui_configured)
    slide_ports(presentation)
    slide_targets(presentation)
    slide_search_size(presentation)
    slide_cascade(presentation, gui_result)
    slide_optimization(presentation)
    slide_results(presentation, decision_plot)
    slide_exports(presentation)
    slide_troubleshooting(presentation)
    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
